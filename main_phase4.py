from pptx import Presentation
from src.ai_processor import AIProcessor
from src.slide_generator import SlideGenerator
import os
import time

def extract_raw_text_per_slide(pptx_path):
    """
    單純提取每一頁的所有文字，不分標題或內文，全部丟進一個 List。
    這是為了模擬「非結構化內容」的情境，讓 AI 來重組。
    """
    print(f"Extracting raw text from: {pptx_path}")
    prs = Presentation(pptx_path)
    slides_raw_data = []

    for i, slide in enumerate(prs.slides):
        text_list = []
        # 遍歷所有 Shape，只要有文字就抓出來
        # 為了讓 AI 更好判斷，我們可以依據 top 位置排序
        shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes.append(shape)

        # 依據垂直位置排序 (由上而下)
        shapes.sort(key=lambda s: s.top)

        for shape in shapes:
            text_list.append(shape.text.strip())

        slides_raw_data.append(text_list)

    return slides_raw_data

def main():
    input_pptx = "input.pptx"
    ref_pptx = "ref.pptx"
    output_pptx = "output_phase4_ai.pptx"

    if not os.path.exists(input_pptx):
        print("Error: Input file not found.")
        return

    # 1. 初始化 AI 處理器
    ai_processor = AIProcessor()

    # 2. 提取原始雜亂文字
    raw_slides = extract_raw_text_per_slide(input_pptx)

    # 3. 利用 AI 重組結構
    structured_slides = []
    print("--- Starting AI Analysis ---")
    for i, raw_text_list in enumerate(raw_slides):
        print(f"Processing Slide {i+1} with AI...")
        if not raw_text_list:
            structured_slides.append({"title": "", "body": []})
            continue

        structured_data = ai_processor.structure_slide_content(raw_text_list)
        print(f"  -> AI Identified Title: {structured_data.get('title')}")
        structured_slides.append(structured_data)

        # Rate Limit Handling: Free tier allows ~15 RPM (1 request every 4 seconds)
        # Adding a small delay to avoid 429 errors
        time.sleep(4)
    print("--- AI Analysis Complete ---")

    # 4. 生成投影片 (使用 Phase 3 的智慧映射引擎)
    # 轉換資料格式以符合 SlideGenerator 的需求 (String -> Complex Run Object)
    final_slides_data = []
    for slide in structured_slides:
        formatted_body = []
        for line in slide.get("body", []):
            # 每一行視為一個段落，每個段落包含一個 Run
            # 預設不加粗、不斜體
            formatted_body.append([
                {"text": line, "bold": False, "italic": False}
            ])

        final_slides_data.append({
            "title": slide.get("title", ""),
            "body": formatted_body
        })

    generator = SlideGenerator(ref_pptx)
    generator.generate(final_slides_data, output_pptx)

    print(f"Success! AI-enhanced output saved to {output_pptx}")

if __name__ == "__main__":
    main()
