from pptx import Presentation

class ContentExtractor:
    def extract(self, pptx_path):
        print(f"Extracting content from: {pptx_path}")
        prs = Presentation(pptx_path)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "title": "",
                "body": [] # List of paragraphs, where each paragraph is a list of runs
            }

            # 1. Identify Title
            # Use the built-in title shape identification
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # 2. Identify Body Content
            # We iterate through all shapes.
            # If it's a text shape and NOT the title, we treat it as body.
            # We sort them by vertical position (top to bottom) to maintain reading order roughly.
            text_shapes = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text_frame") and shape.text.strip():
                    text_shapes.append(shape)

            # Sort by top position
            text_shapes.sort(key=lambda s: s.top)

            for shape in text_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    para_data = []
                    for run in paragraph.runs:
                        run_data = {
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic
                        }
                        para_data.append(run_data)
                    if para_data:
                        slide_data["body"].append(para_data)

            slides_data.append(slide_data)
            print(f"  - Slide {i+1} extracted: Title='{slide_data['title']}', Body Paragraphs={len(slide_data['body'])}")

        return slides_data
