from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

class SlideGenerator:
    def __init__(self, template_path):
        print(f"Loading template from: {template_path}")
        self.template_path = template_path
        self.prs = Presentation(template_path)

        # Analyze layouts to find indices
        self.layout_map = {}
        for i, layout in enumerate(self.prs.slide_layouts):
            self.layout_map[layout.name] = i

        # Define mappings
        self.TITLE_LAYOUT_IDX = self.layout_map.get("TITLE", 0)
        self.SECTION_HEADER_LAYOUT_IDX = self.layout_map.get("SECTION_HEADER", 1)
        self.CONTENT_LAYOUT_IDX = self.layout_map.get("TITLE_AND_BODY", 2)
        self.TWO_CONTENT_LAYOUT_IDX = self.layout_map.get("TITLE_AND_TWO_COLUMNS", 3)

        print(f"  - Title Layout: {self.TITLE_LAYOUT_IDX}")
        print(f"  - Section Header: {self.SECTION_HEADER_LAYOUT_IDX}")
        print(f"  - Content Layout: {self.CONTENT_LAYOUT_IDX}")
        print(f"  - Two Content Layout: {self.TWO_CONTENT_LAYOUT_IDX}")

    def _select_layout_and_split_content(self, index, slide_data):
        """
        Phase 3: Intelligent Mapping Logic
        Returns: (layout_index, [content_group_1, content_group_2, ...])
        """
        title = slide_data.get("title", "")
        body = slide_data.get("body", [])

        # Rule 1: First slide -> Title Slide
        if index == 0:
            # Title slide usually has subtitle in the second placeholder
            return self.TITLE_LAYOUT_IDX, [body]

        # Rule 2: Only Title (no body) -> Section Header
        if not body:
            return self.SECTION_HEADER_LAYOUT_IDX, []

        # Rule 3: Title + Many items -> Two Content
        # Heuristic: If more than 4 paragraphs (lowered threshold for demo), split into two columns
        if len(body) > 4:
            mid = (len(body) + 1) // 2
            col1 = body[:mid]
            col2 = body[mid:]
            print(f"    [Smart Layout] Slide {index+1}: Detected {len(body)} paragraphs -> Using Two Columns")
            return self.TWO_CONTENT_LAYOUT_IDX, [col1, col2]

        # Rule 4: Default -> Title and Content
        return self.CONTENT_LAYOUT_IDX, [body]

    def generate(self, slides_data, output_path):
        print(f"Generating presentation to: {output_path}")

        # Store original slide count to remove them later
        original_slide_count = len(self.prs.slides)

        for i, slide_data in enumerate(slides_data):
            # Phase 3: Use intelligent selection
            layout_idx, content_groups = self._select_layout_and_split_content(i, slide_data)

            slide_layout = self.prs.slide_layouts[layout_idx]
            slide = self.prs.slides.add_slide(slide_layout)

            # 1. Set Title
            if slide_data["title"] and slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]

            # 2. Set Body (handling multiple placeholders for columns)
            # Find all body placeholders (idx > 0)
            body_placeholders = [shape for shape in slide.placeholders if shape.placeholder_format.idx > 0]
            # Sort by idx to ensure we fill left to right (usually idx 1 is left, idx 2 is right)
            body_placeholders.sort(key=lambda x: x.placeholder_format.idx)

            for ph_idx, content_group in enumerate(content_groups):
                if ph_idx < len(body_placeholders):
                    ph = body_placeholders[ph_idx]
                    tf = ph.text_frame
                    tf.clear()

                    # Phase 3: Overflow Handling - Auto Fit
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                    for p_idx, para_data in enumerate(content_group):
                        p = tf.add_paragraph()
                        # Handle the empty first paragraph issue if needed,
                        # but add_paragraph() after clear() usually works fine if we don't mind one empty line or we access paragraphs[0]
                        # Let's stick to add_paragraph for simplicity, or reuse the first one if it's empty.
                        if p_idx == 0 and len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
                            p = tf.paragraphs[0]

                        for run_data in para_data:
                            run = p.add_run()
                            run.text = run_data["text"]
                            if run_data["bold"]:
                                run.font.bold = True
                            if run_data["italic"]:
                                run.font.italic = True

        # Remove original slides after generation to avoid corruption
        if original_slide_count > 0:
            xml_slides = self.prs.slides._sldIdLst
            for _ in range(original_slide_count):
                if len(xml_slides) > 0:
                    xml_slides.remove(xml_slides[0])
            print(f"Removed {original_slide_count} original template slides.")

        self.prs.save(output_path)
        print("Generation complete.")

    def _clear_existing_slides(self):
        # Helper to remove all slides from the template presentation
        # Accessing the private _sldIdLst to clear slides
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)
        for s in slides:
            xml_slides.remove(s)
        print(f"Cleared {len(slides)} existing slides from template.")
