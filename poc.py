import os
from pptx import Presentation

def analyze_content(content_path):
    print(f"--- Task 1.1: Analyzing Content from '{content_path}' ---")
    if not os.path.exists(content_path):
        print(f"Error: File '{content_path}' not found.")
        return

    prs = Presentation(content_path)
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  - Text: {shape.text}")
    print("-" * 30)

def analyze_template(template_path):
    print(f"--- Task 1.2: Analyzing Template from '{template_path}' ---")
    if not os.path.exists(template_path):
        print(f"Error: File '{template_path}' not found.")
        return

    prs = Presentation(template_path)
    print("Available Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")
    print("-" * 30)

def generate_demo(template_path, output_path):
    print(f"--- Task 1.3: Generating Demo to '{output_path}' ---")
    if not os.path.exists(template_path):
        print(f"Error: File '{template_path}' not found.")
        return

    prs = Presentation(template_path)

    # Use the first layout (usually Title Slide) or the second (Title and Content)
    # Let's try to find a layout named "Title Slide" or just use index 0
    layout_to_use = prs.slide_layouts[0]
    print(f"Using Layout: {layout_to_use.name}")

    slide = prs.slides.add_slide(layout_to_use)

    # Fill placeholders
    # Usually placeholder 0 is Title, 1 is Subtitle/Content
    try:
        title = slide.shapes.title
        title.text = "Hello World"
        print("Set Title: Hello World")
    except AttributeError:
        print("No Title placeholder found.")

    # Try to set subtitle or content if available
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = "This is a generated slide."
        print("Set Body: This is a generated slide.")

    prs.save(output_path)
    print(f"Saved output to '{output_path}'")
    print("-" * 30)

if __name__ == "__main__":
    input_pptx = "input.pptx"
    ref_pptx = "ref.pptx"
    output_pptx = "output.pptx"

    analyze_content(input_pptx)
    analyze_template(ref_pptx)
    generate_demo(ref_pptx, output_pptx)
